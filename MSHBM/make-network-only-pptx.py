import os
import sys
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


# Set up argument parser for command-line inputs
parser = argparse.ArgumentParser(description='Make pptx of networks given screenshots.')
parser.add_argument('subid', type=str, help='Subject ID')
parser.add_argument('sess', type=str, help='Session ID')
parser.add_argument('base', type=str, help='Base directory path')
parser.add_argument('codedir', type=str, help='Code directory')
args = parser.parse_args()

subid = args.subid
sess = args.sess
base = args.base
codedir = args.codedir

localroot = '{}/{}/{}_{}/NETWORKS'.format(base, subid, subid, sess)

NETWORKS=['DN-A', "DN-B", "LANG", "FPN-A", "FPN-B", "CGOP", "SALPMN"]

# Define the height for the images
image_height = Inches(1.02)
image_height2 = Inches(1.97)
image_height3 = Inches(2.17)


# Load PowerPoint template
template_path = os.path.join(codedir, "MSHBM/Networks_Template_Blank.pptx")
prs = Presentation(template_path)


# Dynamically create the replacements dictionary
replacements = {
    "SUBID": subid
}


# Access the slide you want to update (e.g., first slide)
slide = prs.slides[0]  

row1 = 0.45
row2 = 1.55
row3 = 2.6
row4 = 4.75
row5 = 5.9
row6 = 6.9


left_position = Inches(1.13)  # Adjust the multiplier as needed to position the images appropriately

top_position = Inches(row1)  
# Add image to slide
image = os.path.join(localroot, "MSHBM_outputs", "L_Lateral.png")
slide.shapes.add_picture(image, left_position, top_position, height=image_height2)
# Add image to slide
top_position = Inches(row3)  
image = os.path.join(localroot, "MSHBM_outputs", "L_Medial.png")
slide.shapes.add_picture(image, left_position, top_position, height=image_height2)
# Add image to slide
top_position = Inches(row4) 
image = os.path.join(localroot, "MSHBM_outputs", "R_Lateral.png")
slide.shapes.add_picture(image, left_position, top_position, height=image_height2)
# Add image to slide
top_position = Inches(row6)  
image = os.path.join(localroot, "MSHBM_outputs", "R_Medial.png")
slide.shapes.add_picture(image, left_position, top_position, height=image_height2)



left_position1 = Inches(-1.13)  # Adjust the multiplier as needed to position the images appropriately
left_position2 = Inches(2.87)  # Adjust the multiplier as needed to position the images appropriately

top_position = Inches(row2)  
# Add image to slide
image = os.path.join(localroot, "MSHBM_outputs", "L_Anterior.png")
slide.shapes.add_picture(image, left_position1, top_position, height=image_height3)
# Add image to slide
image = os.path.join(localroot, "MSHBM_outputs", "L_Posterior.png")
slide.shapes.add_picture(image, left_position2, top_position, height=image_height3)


top_position = Inches(row5)  
# Add image to slide
image = os.path.join(localroot, "MSHBM_outputs", "R_Posterior.png")
slide.shapes.add_picture(image, left_position1, top_position, height=image_height3)
# Add image to slide
image = os.path.join(localroot, "MSHBM_outputs", "R_Anterior.png")
slide.shapes.add_picture(image, left_position2, top_position, height=image_height3)


# ### All slides
# Update labels (text boxes)
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                # Iterate through each run (segment of text with the same formatting) in the paragraph
                for run in paragraph.runs:
                    if "SUBID" in run.text:
                        # Replace "SUBID" with the new text, while keeping the rest of the run's text
                        run.text = run.text.replace("SUBID", subid) 



# Save the updated presentation
out=os.path.join(localroot, '{}_Networks.pptx'.format(subid))
prs.save(out)


print(f"{out} saved.")

